#!/usr/bin/env python3
"""
extract_master_styles.py — Part of the presentation-helper skill.

COMPREHENSIVE PowerPoint master slide and layout style extractor.

Extracts ALL styling and structural information from any PowerPoint file (.pptx)
to produce a complete JSON style reference. This reference enables precise
replication of brand identity in generated presentations.

=== FULL EXTRACTION SCOPE ===

THEME (Complete):
  - Color scheme (dk1, dk2, lt1, lt2, accent1-6, hlink, folHlink)
  - Font scheme (major/minor fonts with scripts)
  - Format scheme (fill styles, line styles, effect styles)
  - Effect scheme (shadows, reflections, 3D effects)
  - Background fill styles
  - Theme name and metadata

MASTER SLIDE (Complete):
  - Background (solid, gradient, picture, pattern fills)
  - All placeholders with full properties
  - Text styles at all 9 levels (font, size, bullets, spacing)
  - Title style definitions
  - Body style definitions
  - Other style definitions (date, footer, slide number)
  - Color map and overrides
  - Decorative shapes with z-order
  - Embedded images/logos with metadata

SLIDE LAYOUTS (Complete per layout):
  - Background (inherit or override)
  - All placeholders (type, idx, position, styles)
  - Decorative shapes
  - Color map overrides
  - Master reference
  - Show/hide master shapes flag

SHAPES (Comprehensive):
  - Position, size (EMU + inches)
  - Shape type (autoshape, picture, group, connector, etc.)
  - Autoshape geometry type (rectangle, oval, arrow, etc.)
  - Fill (solid, gradient, picture, pattern, none)
  - Line/stroke (color, width, dash, compound, cap, join)
  - Shadow effects (inner, outer, perspective)
  - Glow effects
  - Soft edges
  - 3D format and rotation
  - Transparency/opacity
  - Shape locking/protection
  - Z-order

TEXT FRAMES (Complete):
  - Margins (left, right, top, bottom)
  - Vertical anchor
  - Word wrap
  - Auto-size
  - Columns
  - Rotation
  - All paragraphs with:
    - Alignment, spacing (before/after/line)
    - Indentation (left, right, first line, hanging)
    - Bullet/numbering (type, char, font, color, size)
    - Tab stops
    - All runs with:
      - Font (name, size, bold, italic, underline styles)
      - Color (solid, theme reference)
      - Effects (strikethrough, caps, subscript, superscript)
      - Hyperlinks
      - Spacing (kerning, tracking)

PLACEHOLDERS (All 18+ Types):
  - TITLE, CENTER_TITLE, SUBTITLE, BODY, OBJECT
  - CHART, TABLE, SMART_ART, MEDIA_CLIP
  - PICTURE, CLIP_ART, ORG_CHART
  - DATE, FOOTER, SLIDE_NUMBER, HEADER
  - BITMAP, MIXED, VERTICAL variants

OTHER ELEMENTS:
  - Tables (row/column structure, cell styles, borders)
  - Charts (type, embedded data reference)
  - SmartArt (type detection)
  - Grouped shapes (recursive extraction)
  - Connectors (start/end connections)
  - OLE objects
  - Media (video, audio)
  - Notes master and handout master

ACTUAL SLIDES:
  - Full extraction per slide for reference
  - Classification by layout type
  - Notes content

VALIDATION & REPORTING:
  - Extraction warnings for inaccessible elements
  - Coverage statistics
  - Completeness checks

Requirements:
  pip install python-pptx lxml

Usage:
  python extract_master_styles.py <path_to.pptx> [--output <output.json>] [--verbose]

Examples:
  python extract_master_styles.py presentation.pptx
  python extract_master_styles.py deck.pptx -o styles.json --verbose
"""

import sys
import json
import argparse
from pathlib import Path
from collections import defaultdict
from typing import Any, Dict, List, Optional, Union
import warnings

try:
    from pptx import Presentation
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
    from pptx.enum.dml import MSO_LINE_DASH_STYLE, MSO_THEME_COLOR
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
    from lxml import etree
except ImportError:
    print("Installing required packages...")
    import subprocess
    subprocess.check_call(
        [sys.executable, "-m", "pip", "install", "python-pptx", "lxml", "--break-system-packages", "-q"]
    )
    from pptx import Presentation
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
    from pptx.enum.dml import MSO_LINE_DASH_STYLE, MSO_THEME_COLOR
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
    from lxml import etree


# ---------------------------------------------------------------------------
# Global extraction warnings/errors collector
# ---------------------------------------------------------------------------
_extraction_warnings: List[str] = []
_extraction_stats: Dict[str, int] = defaultdict(int)


# ---------------------------------------------------------------------------
# Global extraction warnings/errors collector
# ---------------------------------------------------------------------------
_extraction_warnings: List[str] = []
_extraction_stats: Dict[str, int] = defaultdict(int)


def warn(msg: str):
    """Record an extraction warning."""
    _extraction_warnings.append(msg)


def stat(key: str, increment: int = 1):
    """Increment an extraction statistic."""
    _extraction_stats[key] += increment


# ---------------------------------------------------------------------------
# XML Namespace constants for direct XML access
# ---------------------------------------------------------------------------
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

NSMAP = {
    'a': NS_A,
    'p': NS_P,
    'r': NS_R,
}


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def emu_to_inches(emu: Optional[int]) -> Optional[float]:
    """Convert EMU to inches."""
    if emu is None:
        return None
    return round(emu / 914400, 4)


def emu_to_pt(emu: Optional[int]) -> Optional[float]:
    """Convert EMU to points."""
    if emu is None:
        return None
    return round(emu / 12700, 2)


def rgb_to_hex(rgb) -> Optional[str]:
    """Convert RGBColor to hex string."""
    if rgb is None:
        return None
    try:
        return f"#{rgb.red:02X}{rgb.green:02X}{rgb.blue:02X}"
    except Exception:
        return None


def safe_enum_name(enum_val) -> Optional[str]:
    """Safely get enum name."""
    if enum_val is None:
        return None
    try:
        return enum_val.name if hasattr(enum_val, 'name') else str(enum_val)
    except Exception:
        return str(enum_val)


def safe_int(val) -> Optional[int]:
    """Safely convert to int."""
    if val is None:
        return None
    try:
        return int(val)
    except Exception:
        return None


def get_color_info(color_format) -> Dict[str, Any]:
    """
    Extract complete color information including theme references.
    """
    if color_format is None:
        return {}
    
    info = {}
    try:
        # Color type
        if color_format.type is not None:
            info["type"] = safe_enum_name(color_format.type)
        
        # RGB value
        try:
            rgb = color_format.rgb
            if rgb:
                info["rgb_hex"] = rgb_to_hex(rgb)
        except Exception:
            pass
        
        # Theme color reference
        try:
            theme_color = color_format.theme_color
            if theme_color is not None:
                info["theme_color"] = safe_enum_name(theme_color)
        except Exception:
            pass
        
        # Brightness/luminance modification
        try:
            brightness = color_format.brightness
            if brightness is not None and brightness != 0:
                info["brightness"] = round(brightness, 4)
        except Exception:
            pass
        
    except Exception as e:
        warn(f"Color extraction error: {e}")
    
    return info


def placeholder_type_name(ph_type) -> str:
    """Convert placeholder type to readable name."""
    if ph_type is None:
        return "UNKNOWN"
    try:
        return PP_PLACEHOLDER(ph_type).name
    except Exception:
        return f"TYPE_{ph_type}"


# All known placeholder types for validation
ALL_PLACEHOLDER_TYPES = {
    1: "TITLE",
    2: "BODY",
    3: "CENTER_TITLE",
    4: "SUBTITLE",
    5: "DATE",
    6: "SLIDE_NUMBER",
    7: "FOOTER",
    8: "HEADER",
    9: "OBJECT",
    10: "CHART",
    11: "TABLE",
    12: "CLIP_ART",
    13: "ORG_CHART",
    14: "MEDIA_CLIP",
    15: "BITMAP",
    16: "MIXED",
    17: "BLANK",
    18: "PICTURE",
    19: "VERTICAL_TITLE",
    20: "VERTICAL_BODY",
    21: "VERTICAL_OBJECT",
    22: "SMART_ART",
}


# ---------------------------------------------------------------------------
# Font extraction (comprehensive)
# ---------------------------------------------------------------------------

def extract_font_info(font) -> Dict[str, Any]:
    """
    Extract COMPLETE font information.
    """
    if font is None:
        return {}
    
    info = {}
    
    # Basic font properties
    font_attrs = [
        ("name", "name"),
        ("bold", "bold"),
        ("italic", "italic"),
    ]
    
    for attr, key in font_attrs:
        try:
            val = getattr(font, attr, None)
            if val is not None:
                info[key] = val
        except Exception:
            pass
    
    # Underline (can be boolean or underline type enum)
    try:
        underline = font.underline
        if underline is not None:
            if isinstance(underline, bool):
                info["underline"] = underline
            else:
                info["underline"] = safe_enum_name(underline)
    except Exception:
        pass
    
    # Font size
    try:
        if font.size:
            info["size_pt"] = round(font.size.pt, 1)
            info["size_emu"] = font.size.emu
    except Exception:
        pass
    
    # Color
    try:
        color_info = get_color_info(font.color)
        if color_info:
            info["color"] = color_info
    except Exception:
        pass
    
    # All caps / small caps
    try:
        if font.all_caps is not None:
            info["all_caps"] = font.all_caps
    except Exception:
        pass
    
    # Strikethrough
    try:
        if hasattr(font, 'strikethrough') and font.strikethrough is not None:
            info["strikethrough"] = safe_enum_name(font.strikethrough)
    except Exception:
        pass
    
    # Subscript / superscript (baseline offset)
    try:
        if hasattr(font, 'baseline'):
            baseline = font.baseline
            if baseline is not None and baseline != 0:
                info["baseline_offset"] = baseline  # negative = subscript, positive = superscript
    except Exception:
        pass
    
    # Character spacing / kerning
    try:
        if hasattr(font, '_element'):
            elem = font._element
            cs = elem.get('{%s}spc' % NS_A)
            if cs:
                info["char_spacing"] = int(cs)
    except Exception:
        pass
    
    # Latin / EA / CS font fallbacks
    try:
        if hasattr(font, '_element'):
            elem = font._element
            for font_type in ['latin', 'ea', 'cs']:
                font_elem = elem.find(f'.//{{{NS_A}}}{font_type}')
                if font_elem is not None:
                    typeface = font_elem.get('typeface')
                    if typeface:
                        info[f'{font_type}_font'] = typeface
    except Exception:
        pass
    
    return info


# ---------------------------------------------------------------------------
# Bullet / numbering extraction
# ---------------------------------------------------------------------------

def extract_bullet_info(paragraph) -> Dict[str, Any]:
    """
    Extract bullet/numbering information from paragraph.
    """
    info = {}
    
    try:
        pPr = None
        if hasattr(paragraph, '_element') and hasattr(paragraph._element, 'pPr'):
            pPr = paragraph._element.pPr
        elif hasattr(paragraph, '_p') and hasattr(paragraph._p, 'pPr'):
            pPr = paragraph._p.pPr
        
        if pPr is None:
            return info
        
        # Check for bullet character
        buChar = pPr.find(f'.//{{{NS_A}}}buChar')
        if buChar is not None:
            info["bullet_type"] = "character"
            info["bullet_char"] = buChar.get('char')
        
        # Check for auto-numbered bullets
        buAutoNum = pPr.find(f'.//{{{NS_A}}}buAutoNum')
        if buAutoNum is not None:
            info["bullet_type"] = "auto_number"
            info["numbering_type"] = buAutoNum.get('type')
            start = buAutoNum.get('startAt')
            if start:
                info["start_at"] = int(start)
        
        # Check for picture bullets
        buBlip = pPr.find(f'.//{{{NS_A}}}buBlip')
        if buBlip is not None:
            info["bullet_type"] = "picture"
        
        # Bullet color
        buClr = pPr.find(f'.//{{{NS_A}}}buClr')
        if buClr is not None:
            srgb = buClr.find(f'.//{{{NS_A}}}srgbClr')
            if srgb is not None:
                info["bullet_color"] = f"#{srgb.get('val', '000000').upper()}"
        
        # Bullet size (percent or points)
        buSzPct = pPr.find(f'.//{{{NS_A}}}buSzPct')
        if buSzPct is not None:
            val = buSzPct.get('val')
            if val:
                info["bullet_size_pct"] = int(val) / 1000  # stored as 1000ths
        
        buSzPts = pPr.find(f'.//{{{NS_A}}}buSzPts')
        if buSzPts is not None:
            val = buSzPts.get('val')
            if val:
                info["bullet_size_pt"] = int(val) / 100  # stored as 100ths
        
        # Bullet font
        buFont = pPr.find(f'.//{{{NS_A}}}buFont')
        if buFont is not None:
            info["bullet_font"] = buFont.get('typeface')
        
        # No bullet
        buNone = pPr.find(f'.//{{{NS_A}}}buNone')
        if buNone is not None:
            info["bullet_type"] = "none"
        
    except Exception as e:
        warn(f"Bullet extraction error: {e}")
    
    return info


# ---------------------------------------------------------------------------
# Paragraph extraction (comprehensive)
# ---------------------------------------------------------------------------

def extract_paragraph_info(paragraph, include_text: bool = False) -> Dict[str, Any]:
    """
    Extract COMPLETE paragraph information.
    """
    info = {}
    
    # Check if paragraph has paragraph_format (not all _Paragraph types do)
    if hasattr(paragraph, 'paragraph_format'):
        try:
            pf = paragraph.paragraph_format
            
            # Alignment
            if pf.alignment is not None:
                info["alignment"] = safe_enum_name(pf.alignment)
            
            # Spacing before/after
            if pf.space_before is not None:
                try:
                    info["space_before_pt"] = round(pf.space_before.pt, 1)
                except Exception:
                    pass
            
            if pf.space_after is not None:
                try:
                    info["space_after_pt"] = round(pf.space_after.pt, 1)
                except Exception:
                    pass
            
            # Line spacing
            if pf.line_spacing is not None:
                info["line_spacing"] = pf.line_spacing
                if hasattr(pf, 'line_spacing_rule') and pf.line_spacing_rule is not None:
                    info["line_spacing_rule"] = safe_enum_name(pf.line_spacing_rule)
            
            # Indentation
            try:
                if pf.left_margin is not None:
                    info["left_margin_in"] = emu_to_inches(pf.left_margin)
                if pf.right_margin is not None:
                    info["right_margin_in"] = emu_to_inches(pf.right_margin)
                if pf.first_line_indent is not None:
                    info["first_line_indent_in"] = emu_to_inches(pf.first_line_indent)
            except Exception:
                pass
        
        except Exception:
            pass
    
    # Level - try different approaches
    try:
        if hasattr(paragraph, 'level'):
            info["level"] = paragraph.level
    except Exception:
        pass

    # Bullet/numbering
    bullet_info = extract_bullet_info(paragraph)
    if bullet_info:
        info["bullet"] = bullet_info

    # Font info from all runs (merge unique properties)
    merged_font = {}
    runs_info = []
    
    for run in paragraph.runs:
        run_info = {}
        rf = extract_font_info(run.font)
        if rf:
            run_info["font"] = rf
            for k, v in rf.items():
                if k not in merged_font and v is not None:
                    merged_font[k] = v
        
        # Hyperlink detection
        try:
            if hasattr(run, '_r') and hasattr(run._r, 'hlinkClick'):
                hlink = run._r.hlinkClick
                if hlink is not None:
                    run_info["has_hyperlink"] = True
                    rId = hlink.get('{%s}id' % NS_R)
                    if rId:
                        run_info["hyperlink_rId"] = rId
        except Exception:
            pass
        
        if include_text:
            try:
                text = run.text
                if text:
                    run_info["text"] = text[:50]
            except Exception:
                pass
        
        if run_info:
            runs_info.append(run_info)
    
    if merged_font:
        info["font"] = merged_font
    
    if runs_info and len(runs_info) > 1:
        # Only include runs detail if there are multiple with different formatting
        info["runs_detail"] = runs_info

    if include_text:
        try:
            text = paragraph.text.strip()
            if text:
                info["sample_text"] = text[:100]
        except Exception:
            pass

    return {k: v for k, v in info.items() if v is not None}


# ---------------------------------------------------------------------------
# Fill extraction (comprehensive)
# ---------------------------------------------------------------------------

def extract_fill_info(fill) -> Dict[str, Any]:
    """
    Extract COMPLETE fill information including all fill types.
    """
    if fill is None:
        return {}
    
    info = {}
    
    try:
        fill_type = str(fill.type) if fill.type is not None else None
        info["type"] = fill_type
        
        if not fill_type:
            return info
        
        upper = fill_type.upper()
        
        if "SOLID" in upper:
            # Solid fill
            try:
                info["color"] = get_color_info(fill.fore_color)
            except Exception:
                pass
            
            # Transparency
            try:
                if hasattr(fill, 'transparency') and fill.transparency is not None:
                    info["transparency"] = round(fill.transparency, 4)
            except Exception:
                pass
                
        elif "GRADIENT" in upper:
            # Gradient fill
            try:
                stops = []
                for stop in fill.gradient_stops:
                    stop_info = {
                        "position": round(stop.position, 4),
                        "color": get_color_info(stop.color)
                    }
                    stops.append(stop_info)
                info["gradient_stops"] = stops
            except Exception as e:
                warn(f"Gradient stops extraction error: {e}")
            
            try:
                if fill.gradient_angle is not None:
                    info["gradient_angle"] = fill.gradient_angle
            except Exception:
                pass
            
            # Gradient type (linear, radial, etc.)
            try:
                if hasattr(fill, 'gradient_type') and fill.gradient_type is not None:
                    info["gradient_type"] = safe_enum_name(fill.gradient_type)
            except Exception:
                pass
                
        elif "PATTERN" in upper:
            # Pattern fill
            info["fill_category"] = "pattern"
            try:
                info["fore_color"] = get_color_info(fill.fore_color)
                info["back_color"] = get_color_info(fill.back_color)
            except Exception:
                pass
            try:
                if hasattr(fill, 'pattern') and fill.pattern is not None:
                    info["pattern_type"] = safe_enum_name(fill.pattern)
            except Exception:
                pass
                
        elif "PICTURE" in upper or "TEXTURE" in upper:
            # Picture/texture fill
            info["fill_category"] = "picture_or_texture"
            info["note"] = "background image or texture present"
            
            # Try to get image metadata
            try:
                if hasattr(fill, '_element'):
                    blipFill = fill._element.find(f'.//{{{NS_A}}}blipFill')
                    if blipFill is not None:
                        blip = blipFill.find(f'.//{{{NS_A}}}blip')
                        if blip is not None:
                            embed = blip.get('{%s}embed' % NS_R)
                            if embed:
                                info["image_rId"] = embed
            except Exception:
                pass
            
            # Tile/stretch settings
            try:
                if hasattr(fill, '_element'):
                    tile = fill._element.find(f'.//{{{NS_A}}}tile')
                    if tile is not None:
                        info["tile_mode"] = True
                        info["tile_sx"] = tile.get('sx')
                        info["tile_sy"] = tile.get('sy')
                    
                    stretch = fill._element.find(f'.//{{{NS_A}}}stretch')
                    if stretch is not None:
                        info["stretch_mode"] = True
            except Exception:
                pass
                
        elif "BACKGROUND" in upper:
            info["inherits_slide_background"] = True
            
        elif "NONE" in upper or "NO_FILL" in upper:
            info["no_fill"] = True
            
    except Exception as e:
        warn(f"Fill extraction error: {e}")
    
    return {k: v for k, v in info.items() if v is not None}


# ---------------------------------------------------------------------------
# Line/stroke extraction (comprehensive)
# ---------------------------------------------------------------------------

def extract_line_info(line) -> Dict[str, Any]:
    """
    Extract COMPLETE line/stroke information.
    """
    if line is None:
        return {}
    
    info = {}
    
    try:
        # Width
        if line.width is not None:
            info["width_pt"] = round(line.width.pt, 2)
            info["width_emu"] = line.width.emu
    except Exception:
        pass
    
    try:
        # Color
        color_info = get_color_info(line.color)
        if color_info:
            info["color"] = color_info
    except Exception:
        pass
    
    try:
        # Dash style
        if line.dash_style is not None:
            info["dash_style"] = safe_enum_name(line.dash_style)
    except Exception:
        pass
    
    try:
        # Compound type (single, double, triple, etc.)
        if hasattr(line, 'compound_type') and line.compound_type is not None:
            info["compound_type"] = safe_enum_name(line.compound_type)
    except Exception:
        pass
    
    try:
        # Cap type (flat, round, square)
        if hasattr(line, 'cap_type') and line.cap_type is not None:
            info["cap_type"] = safe_enum_name(line.cap_type)
    except Exception:
        pass
    
    try:
        # Join type (bevel, miter, round)
        if hasattr(line, 'join_type') and line.join_type is not None:
            info["join_type"] = safe_enum_name(line.join_type)
    except Exception:
        pass
    
    try:
        # Fill (for complex line fills)
        if hasattr(line, 'fill'):
            line_fill = extract_fill_info(line.fill)
            if line_fill and line_fill.get('type'):
                info["fill"] = line_fill
    except Exception:
        pass
    
    return info


# ---------------------------------------------------------------------------
# Effect extraction (shadows, glow, soft edges, 3D)
# ---------------------------------------------------------------------------

def extract_effect_info(shape) -> Dict[str, Any]:
    """
    Extract shape effects (shadow, glow, reflection, soft edges, 3D).
    """
    info = {}
    
    try:
        if not hasattr(shape, '_element'):
            return info
        
        elem = shape._element
        spPr = elem.find(f'.//{{{NS_A}}}spPr')
        if spPr is None:
            return info
        
        effectLst = spPr.find(f'.//{{{NS_A}}}effectLst')
        if effectLst is None:
            return info
        
        # Outer shadow
        outerShdw = effectLst.find(f'.//{{{NS_A}}}outerShdw')
        if outerShdw is not None:
            shadow = {
                "type": "outer",
                "blur_rad": outerShdw.get('blurRad'),
                "dist": outerShdw.get('dist'),
                "dir": outerShdw.get('dir'),
                "sx": outerShdw.get('sx'),
                "sy": outerShdw.get('sy'),
            }
            # Get shadow color
            srgbClr = outerShdw.find(f'.//{{{NS_A}}}srgbClr')
            if srgbClr is not None:
                shadow["color"] = f"#{srgbClr.get('val', '000000').upper()}"
                alpha = srgbClr.find(f'.//{{{NS_A}}}alpha')
                if alpha is not None:
                    shadow["alpha"] = alpha.get('val')
            info["outer_shadow"] = shadow
        
        # Inner shadow
        innerShdw = effectLst.find(f'.//{{{NS_A}}}innerShdw')
        if innerShdw is not None:
            info["inner_shadow"] = {
                "type": "inner",
                "blur_rad": innerShdw.get('blurRad'),
                "dist": innerShdw.get('dist'),
                "dir": innerShdw.get('dir'),
            }
        
        # Glow
        glow = effectLst.find(f'.//{{{NS_A}}}glow')
        if glow is not None:
            glow_info = {"rad": glow.get('rad')}
            srgbClr = glow.find(f'.//{{{NS_A}}}srgbClr')
            if srgbClr is not None:
                glow_info["color"] = f"#{srgbClr.get('val', '000000').upper()}"
            info["glow"] = glow_info
        
        # Soft edges
        softEdge = effectLst.find(f'.//{{{NS_A}}}softEdge')
        if softEdge is not None:
            info["soft_edge"] = {"rad": softEdge.get('rad')}
        
        # Reflection
        reflection = effectLst.find(f'.//{{{NS_A}}}reflection')
        if reflection is not None:
            info["reflection"] = {
                "blur_rad": reflection.get('blurRad'),
                "start_opacity": reflection.get('stA'),
                "end_opacity": reflection.get('endA'),
                "dist": reflection.get('dist'),
                "dir": reflection.get('dir'),
            }
        
        # 3D format
        sp3d = spPr.find(f'.//{{{NS_A}}}sp3d')
        if sp3d is not None:
            info["3d_format"] = {
                "preset_material": sp3d.get('prstMaterial'),
                "extrusion_h": sp3d.get('extrusionH'),
                "contour_w": sp3d.get('contourW'),
            }
        
        # Scene 3D (camera/lighting)
        scene3d = spPr.find(f'.//{{{NS_A}}}scene3d')
        if scene3d is not None:
            camera = scene3d.find(f'.//{{{NS_A}}}camera')
            if camera is not None:
                info["3d_scene"] = {
                    "camera_preset": camera.get('prst'),
                }
        
    except Exception as e:
        warn(f"Effect extraction error: {e}")
    
    return info


# ---------------------------------------------------------------------------
# Shape extraction (comprehensive)
# ---------------------------------------------------------------------------

def extract_shape_info(shape, include_text_samples: bool = False, z_order: int = 0) -> Dict[str, Any]:
    """
    Extract COMPLETE shape information including all properties.
    """
    stat("shapes_extracted")
    
    info = {
        "name": shape.name,
        "shape_type": safe_enum_name(shape.shape_type),
        "shape_type_int": safe_int(shape.shape_type),
        "z_order": z_order,
    }

    # ---- Position and size ----
    try:
        info["position"] = {
            "left_emu": shape.left,
            "top_emu": shape.top,
            "width_emu": shape.width,
            "height_emu": shape.height,
            "left_in": emu_to_inches(shape.left),
            "top_in": emu_to_inches(shape.top),
            "width_in": emu_to_inches(shape.width),
            "height_in": emu_to_inches(shape.height),
        }
    except Exception as e:
        warn(f"Position extraction error for {shape.name}: {e}")

    # ---- Rotation ----
    try:
        if hasattr(shape, 'rotation') and shape.rotation is not None and shape.rotation != 0:
            info["rotation_degrees"] = shape.rotation
    except Exception:
        pass

    # ---- Placeholder metadata ----
    try:
        if shape.is_placeholder:
            ph = shape.placeholder_format
            ph_type_int = safe_int(ph.type)
            info["placeholder"] = {
                "type": placeholder_type_name(ph.type),
                "type_int": ph_type_int,
                "idx": ph.idx,
                "category": _categorize_placeholder(ph_type_int),
            }
            stat(f"placeholder_{placeholder_type_name(ph.type)}")
    except Exception:
        pass

    # ---- Auto-shape geometry type ----
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if hasattr(shape, 'auto_shape_type') and shape.auto_shape_type is not None:
                info["auto_shape_type"] = safe_enum_name(shape.auto_shape_type)
    except Exception:
        pass

    # ---- Image / picture ----
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            info["is_image"] = True
            stat("images")
            try:
                img = shape.image
                info["image_info"] = {
                    "content_type": img.content_type,
                    "size_bytes": len(img.blob),
                    "filename": getattr(img, 'filename', None),
                    "dpi": getattr(img, 'dpi', None),
                }
            except Exception:
                info["image_info"] = {"note": "image blob inaccessible"}
    except Exception:
        pass

    # ---- Picture placeholder ----
    try:
        if shape.is_placeholder:
            ph = shape.placeholder_format
            if ph.type is not None and int(ph.type) == 18:  # PP_PLACEHOLDER.PICTURE
                info["is_picture_placeholder"] = True
                stat("picture_placeholders")
                try:
                    img = shape.image
                    info["image_info"] = {
                        "content_type": img.content_type,
                        "size_bytes": len(img.blob),
                        "note": "picture placeholder with image"
                    }
                except Exception:
                    info["image_info"] = {"note": "picture placeholder (empty)"}
    except Exception:
        pass

    # ---- Table ----
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            info["is_table"] = True
            stat("tables")
            table = shape.table
            table_info = {
                "rows": len(table.rows),
                "columns": len(table.columns),
                "first_row_header": getattr(table, 'first_row', None),
                "last_row_banded": getattr(table, 'last_row', None),
                "first_col_banded": getattr(table, 'first_col', None),
                "last_col_banded": getattr(table, 'last_col', None),
                "horz_banding": getattr(table, 'horz_banding', None),
                "vert_banding": getattr(table, 'vert_banding', None),
            }
            
            # Extract cell styles from first row
            if len(table.rows) > 0:
                first_row_cells = []
                for cell in table.rows[0].cells:
                    cell_info = {}
                    try:
                        cell_info["fill"] = extract_fill_info(cell.fill)
                    except Exception:
                        pass
                    try:
                        if cell.text_frame:
                            paras = []
                            for para in cell.text_frame.paragraphs:
                                paras.append(extract_paragraph_info(para, include_text=True))
                            if paras:
                                cell_info["paragraphs"] = paras
                    except Exception:
                        pass
                    first_row_cells.append(cell_info)
                table_info["first_row_cells_sample"] = first_row_cells
            
            info["table_info"] = table_info
    except Exception:
        pass

    # ---- Chart ----
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            info["is_chart"] = True
            stat("charts")
            try:
                chart = shape.chart
                info["chart_info"] = {
                    "chart_type": safe_enum_name(chart.chart_type),
                    "has_legend": chart.has_legend,
                }
            except Exception:
                info["chart_info"] = {"note": "chart detected"}
    except Exception:
        pass

    # ---- SmartArt / Diagram ----
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.SMART_ART:
            info["is_smart_art"] = True
            stat("smartart")
            info["smart_art_info"] = {"note": "SmartArt diagram detected"}
    except Exception:
        pass

    # ---- Group shape (recursive) ----
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            info["is_group"] = True
            stat("groups")
            group_shapes = []
            for idx, child_shape in enumerate(shape.shapes):
                child_info = extract_shape_info(child_shape, include_text_samples, z_order=idx)
                group_shapes.append(child_info)
            info["group_shapes"] = group_shapes
    except Exception:
        pass

    # ---- Connector ----
    try:
        if shape.shape_type in (MSO_SHAPE_TYPE.LINE, 14):  # LINE or CONNECTOR
            info["is_connector"] = True
            stat("connectors")
            # Try to get connector endpoints
            try:
                if hasattr(shape, 'begin_x'):
                    info["connector_info"] = {
                        "begin_x_in": emu_to_inches(shape.begin_x),
                        "begin_y_in": emu_to_inches(shape.begin_y),
                        "end_x_in": emu_to_inches(shape.end_x),
                        "end_y_in": emu_to_inches(shape.end_y),
                    }
            except Exception:
                pass
    except Exception:
        pass

    # ---- OLE Object ----
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
            info["is_ole_object"] = True
            stat("ole_objects")
            info["ole_info"] = {"note": "Embedded OLE object"}
    except Exception:
        pass

    # ---- Media (Video/Audio) ----
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
            info["is_media"] = True
            stat("media")
            info["media_info"] = {"note": "Media clip (video/audio)"}
    except Exception:
        pass

    # ---- Footer / special shape detection ----
    try:
        lower_name = shape.name.lower()
        footer_keywords = ("fußzeile", "fusszeile", "footer", "seitenzahl", 
                          "pagenumber", "page number", "slide number", "date", "datum")
        if any(kw in lower_name for kw in footer_keywords):
            info["is_footer_element"] = True
            if shape.has_text_frame:
                info["footer_text"] = shape.text_frame.text
        
        confidential_keywords = ("confidential", "vertraulich", "internal", "proprietary")
        if any(kw in lower_name for kw in confidential_keywords):
            info["is_confidential_label"] = True
        elif shape.has_text_frame:
            text_upper = shape.text_frame.text.upper()
            if any(kw.upper() in text_upper for kw in confidential_keywords):
                info["is_confidential_label"] = True
    except Exception:
        pass

    # ---- Fill ----
    try:
        fill_info = extract_fill_info(shape.fill)
        if fill_info:
            info["fill"] = fill_info
    except Exception:
        pass

    # ---- Line / border ----
    try:
        line_info = extract_line_info(shape.line)
        if line_info:
            info["line"] = line_info
    except Exception:
        pass

    # ---- Effects (shadow, glow, etc.) ----
    try:
        effects = extract_effect_info(shape)
        if effects:
            info["effects"] = effects
    except Exception:
        pass

    # ---- Shape locking ----
    try:
        if hasattr(shape, '_element'):
            cNvSpPr = shape._element.find(f'.//{{{NS_P}}}cNvSpPr')
            if cNvSpPr is None:
                cNvSpPr = shape._element.find(f'.//{{{NS_A}}}cNvSpPr')
            if cNvSpPr is not None:
                locks = cNvSpPr.find(f'.//{{{NS_A}}}spLocks')
                if locks is not None:
                    lock_info = {}
                    for attr in ['noMove', 'noResize', 'noRot', 'noChangeAspect', 
                                'noEditPoints', 'noAdjustHandles', 'noChangeArrowheads',
                                'noChangeShapeType', 'noTextEdit', 'noGrp']:
                        val = locks.get(attr)
                        if val == '1' or val == 'true':
                            lock_info[attr] = True
                    if lock_info:
                        info["locks"] = lock_info
    except Exception:
        pass

    # ---- Text frame ----
    try:
        if shape.has_text_frame:
            tf = shape.text_frame
            tf_info = {}
            
            # Auto-size
            if tf.auto_size is not None:
                tf_info["auto_size"] = safe_enum_name(tf.auto_size)
            
            # Word wrap
            if tf.word_wrap is not None:
                tf_info["word_wrap"] = tf.word_wrap
            
            # Margins
            if tf.margin_left is not None:
                tf_info["margin_left_in"] = emu_to_inches(tf.margin_left)
            if tf.margin_right is not None:
                tf_info["margin_right_in"] = emu_to_inches(tf.margin_right)
            if tf.margin_top is not None:
                tf_info["margin_top_in"] = emu_to_inches(tf.margin_top)
            if tf.margin_bottom is not None:
                tf_info["margin_bottom_in"] = emu_to_inches(tf.margin_bottom)
            
            # Vertical anchor
            if tf.vertical_anchor:
                tf_info["vertical_anchor"] = safe_enum_name(tf.vertical_anchor)
            
            # Text body properties from XML
            try:
                if hasattr(tf, '_txBody'):
                    bodyPr = tf._txBody.find(f'.//{{{NS_A}}}bodyPr')
                    if bodyPr is not None:
                        # Rotation
                        rot = bodyPr.get('rot')
                        if rot:
                            tf_info["text_rotation"] = int(rot) / 60000  # Convert from EMU angle
                        
                        # Columns
                        num_col = bodyPr.get('numCol')
                        if num_col:
                            tf_info["columns"] = int(num_col)
                        
                        # Upright text
                        upright = bodyPr.get('upright')
                        if upright == '1' or upright == 'true':
                            tf_info["upright"] = True
                        
                        # Vertical text
                        vert = bodyPr.get('vert')
                        if vert:
                            tf_info["text_direction"] = vert
            except Exception:
                pass

            # Paragraphs
            paragraphs = []
            for para in tf.paragraphs:
                para_info = extract_paragraph_info(para, include_text=include_text_samples)
                if para_info:
                    paragraphs.append(para_info)
            if paragraphs:
                tf_info["paragraphs"] = paragraphs
                tf_info["paragraph_count"] = len(paragraphs)

            if tf_info:
                info["text_frame"] = tf_info
    except Exception as e:
        warn(f"Text frame extraction error for {shape.name}: {e}")

    return info


def _categorize_placeholder(ph_type_int: Optional[int]) -> str:
    """Categorize placeholder by type for easier filtering."""
    if ph_type_int is None:
        return "unknown"
    
    title_types = {1, 3, 19}  # TITLE, CENTER_TITLE, VERTICAL_TITLE
    body_types = {2, 20}  # BODY, VERTICAL_BODY
    object_types = {9, 10, 11, 12, 13, 14, 18, 21, 22}  # OBJECT, CHART, TABLE, etc.
    footer_types = {5, 6, 7, 8}  # DATE, SLIDE_NUMBER, FOOTER, HEADER
    
    if ph_type_int in title_types:
        return "title"
    elif ph_type_int in body_types:
        return "body"
    elif ph_type_int == 4:
        return "subtitle"
    elif ph_type_int in object_types:
        return "object"
    elif ph_type_int in footer_types:
        return "footer"
    else:
        return "other"


# ---------------------------------------------------------------------------
# Theme extraction (comprehensive)
# ---------------------------------------------------------------------------

def extract_theme_colors(prs) -> Dict[str, Any]:
    """Extract COMPLETE theme information including colors, fonts, and effects."""
    theme_data = {
        "colors": {},
        "font_scheme": {},
        "format_scheme": {},
        "effect_scheme": {},
        "background_fill_styles": [],
        "theme_name": None,
    }
    
    try:
        for slide_master in prs.slide_masters:
            theme_part = slide_master.part.part_related_by(
                'http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme'
            )
            
            # Get XML tree
            tree = None
            if hasattr(theme_part, '_element'):
                tree = theme_part._element
            elif hasattr(theme_part, 'element'):
                tree = theme_part.element
            elif hasattr(theme_part, 'blob'):
                tree = etree.fromstring(theme_part.blob)
            
            if tree is None:
                warn("Could not access theme XML element")
                return theme_data
            
            # Theme name
            try:
                theme_name = tree.get('name')
                if theme_name:
                    theme_data["theme_name"] = theme_name
            except Exception:
                pass
            
            # ---- Color Scheme ----
            clrScheme = tree.find(f'.//{{{NS_A}}}clrScheme')
            if clrScheme is not None:
                scheme_name = clrScheme.get('name')
                if scheme_name:
                    theme_data["colors"]["scheme_name"] = scheme_name
                
                color_names = ['dk1', 'dk2', 'lt1', 'lt2', 
                               'accent1', 'accent2', 'accent3', 'accent4', 'accent5', 'accent6',
                               'hlink', 'folHlink']
                
                for color_name in color_names:
                    color_elem = clrScheme.find(f'{{{NS_A}}}{color_name}')
                    if color_elem is not None:
                        color_info = _extract_xml_color(color_elem)
                        if color_info:
                            theme_data["colors"][color_name] = color_info
            
            # ---- Font Scheme ----
            fontScheme = tree.find(f'.//{{{NS_A}}}fontScheme')
            if fontScheme is not None:
                scheme_name = fontScheme.get('name')
                if scheme_name:
                    theme_data["font_scheme"]["scheme_name"] = scheme_name
                
                for font_group_tag in ['majorFont', 'minorFont']:
                    font_group = fontScheme.find(f'{{{NS_A}}}{font_group_tag}')
                    if font_group is not None:
                        group_data = {}
                        
                        # Latin font
                        latin = font_group.find(f'{{{NS_A}}}latin')
                        if latin is not None:
                            group_data["latin"] = latin.get('typeface')
                        
                        # East Asian font
                        ea = font_group.find(f'{{{NS_A}}}ea')
                        if ea is not None:
                            group_data["east_asian"] = ea.get('typeface')
                        
                        # Complex script font
                        cs = font_group.find(f'{{{NS_A}}}cs')
                        if cs is not None:
                            group_data["complex_script"] = cs.get('typeface')
                        
                        # Font alternatives for specific locales
                        fonts_by_script = {}
                        for font_elem in font_group.findall(f'{{{NS_A}}}font'):
                            script = font_elem.get('script')
                            typeface = font_elem.get('typeface')
                            if script and typeface:
                                fonts_by_script[script] = typeface
                        if fonts_by_script:
                            group_data["by_script"] = fonts_by_script
                        
                        theme_data["font_scheme"][font_group_tag] = group_data
            
            # ---- Format Scheme (fill, line, effect styles) ----
            fmtScheme = tree.find(f'.//{{{NS_A}}}fmtScheme')
            if fmtScheme is not None:
                scheme_name = fmtScheme.get('name')
                if scheme_name:
                    theme_data["format_scheme"]["scheme_name"] = scheme_name
                
                # Fill styles
                fillStyleLst = fmtScheme.find(f'{{{NS_A}}}fillStyleLst')
                if fillStyleLst is not None:
                    fill_styles = []
                    for idx, fill_elem in enumerate(fillStyleLst):
                        fill_info = _extract_xml_fill(fill_elem)
                        fill_info["index"] = idx + 1
                        fill_styles.append(fill_info)
                    theme_data["format_scheme"]["fill_styles"] = fill_styles
                
                # Line styles
                lnStyleLst = fmtScheme.find(f'{{{NS_A}}}lnStyleLst')
                if lnStyleLst is not None:
                    line_styles = []
                    for idx, ln_elem in enumerate(lnStyleLst):
                        line_info = _extract_xml_line(ln_elem)
                        line_info["index"] = idx + 1
                        line_styles.append(line_info)
                    theme_data["format_scheme"]["line_styles"] = line_styles
                
                # Effect styles
                effectStyleLst = fmtScheme.find(f'{{{NS_A}}}effectStyleLst')
                if effectStyleLst is not None:
                    effect_styles = []
                    for idx, effectStyle in enumerate(effectStyleLst):
                        effect_info = {"index": idx + 1}
                        effectLst = effectStyle.find(f'{{{NS_A}}}effectLst')
                        if effectLst is not None:
                            effects = []
                            for effect in effectLst:
                                effect_name = effect.tag.split('}')[-1]
                                effects.append(effect_name)
                            effect_info["effects"] = effects
                        effect_styles.append(effect_info)
                    theme_data["format_scheme"]["effect_styles"] = effect_styles
                
                # Background fill styles
                bgFillStyleLst = fmtScheme.find(f'{{{NS_A}}}bgFillStyleLst')
                if bgFillStyleLst is not None:
                    bg_fills = []
                    for idx, fill_elem in enumerate(bgFillStyleLst):
                        fill_info = _extract_xml_fill(fill_elem)
                        fill_info["index"] = idx + 1
                        bg_fills.append(fill_info)
                    theme_data["background_fill_styles"] = bg_fills
            
            break  # First master only
            
    except Exception as e:
        warn(f"Theme extraction error: {e}")
    
    return theme_data


def _extract_xml_color(color_elem) -> Optional[Dict[str, Any]]:
    """Extract color from XML element."""
    if color_elem is None:
        return None
    
    info = {}
    
    # Check child color type
    for child in color_elem:
        tag = child.tag.split('}')[-1]
        
        if tag == 'srgbClr':
            val = child.get('val')
            if val:
                info["type"] = "srgb"
                info["hex"] = f"#{val.upper()}"
                # Check for transforms
                _extract_color_transforms(child, info)
                
        elif tag == 'sysClr':
            info["type"] = "system"
            info["system_color"] = child.get('val')
            last_clr = child.get('lastClr')
            if last_clr:
                info["hex"] = f"#{last_clr.upper()}"
            _extract_color_transforms(child, info)
            
        elif tag == 'schemeClr':
            info["type"] = "scheme"
            info["scheme_ref"] = child.get('val')
            _extract_color_transforms(child, info)
    
    return info if info else None


def _extract_color_transforms(color_elem, info: Dict):
    """Extract color transforms (alpha, tint, shade, etc.)."""
    transforms = {}
    
    for child in color_elem:
        tag = child.tag.split('}')[-1]
        val = child.get('val')
        
        if tag == 'alpha' and val:
            transforms["alpha"] = int(val) / 1000
        elif tag == 'tint' and val:
            transforms["tint"] = int(val) / 1000
        elif tag == 'shade' and val:
            transforms["shade"] = int(val) / 1000
        elif tag == 'satMod' and val:
            transforms["saturation_mod"] = int(val) / 1000
        elif tag == 'lumMod' and val:
            transforms["luminance_mod"] = int(val) / 1000
        elif tag == 'lumOff' and val:
            transforms["luminance_off"] = int(val) / 1000
    
    if transforms:
        info["transforms"] = transforms


def _extract_xml_fill(fill_elem) -> Dict[str, Any]:
    """Extract fill info from XML element."""
    info = {}
    tag = fill_elem.tag.split('}')[-1]
    info["fill_type"] = tag
    
    if tag == 'solidFill':
        for color_child in fill_elem:
            color_info = _extract_xml_color_direct(color_child)
            if color_info:
                info["color"] = color_info
                
    elif tag == 'gradFill':
        info["gradient"] = True
        gsLst = fill_elem.find(f'{{{NS_A}}}gsLst')
        if gsLst is not None:
            stops = []
            for gs in gsLst.findall(f'{{{NS_A}}}gs'):
                pos = gs.get('pos')
                stop_info = {"position": int(pos) / 100000 if pos else 0}
                for color_child in gs:
                    color_info = _extract_xml_color_direct(color_child)
                    if color_info:
                        stop_info["color"] = color_info
                stops.append(stop_info)
            info["stops"] = stops
        
        lin = fill_elem.find(f'{{{NS_A}}}lin')
        if lin is not None:
            ang = lin.get('ang')
            if ang:
                info["angle"] = int(ang) / 60000
                
    elif tag == 'pattFill':
        info["pattern"] = fill_elem.get('prst')
        
    elif tag == 'blipFill':
        info["picture_fill"] = True
        
    elif tag == 'noFill':
        info["no_fill"] = True
    
    return info


def _extract_xml_color_direct(color_elem) -> Optional[Dict[str, Any]]:
    """Extract color directly from a color element."""
    if color_elem is None:
        return None
    
    tag = color_elem.tag.split('}')[-1]
    info = {"type": tag}
    
    if tag == 'srgbClr':
        val = color_elem.get('val')
        if val:
            info["hex"] = f"#{val.upper()}"
    elif tag == 'schemeClr':
        info["scheme_ref"] = color_elem.get('val')
    elif tag == 'sysClr':
        info["system"] = color_elem.get('val')
        last = color_elem.get('lastClr')
        if last:
            info["hex"] = f"#{last.upper()}"
    
    _extract_color_transforms(color_elem, info)
    return info


def _extract_xml_line(ln_elem) -> Dict[str, Any]:
    """Extract line info from XML element."""
    info = {}
    
    w = ln_elem.get('w')
    if w:
        info["width_emu"] = int(w)
        info["width_pt"] = round(int(w) / 12700, 2)
    
    cap = ln_elem.get('cap')
    if cap:
        info["cap"] = cap
    
    cmpd = ln_elem.get('cmpd')
    if cmpd:
        info["compound"] = cmpd
    
    # Fill
    solidFill = ln_elem.find(f'{{{NS_A}}}solidFill')
    if solidFill is not None:
        for color_child in solidFill:
            color_info = _extract_xml_color_direct(color_child)
            if color_info:
                info["color"] = color_info
    
    # Dash
    prstDash = ln_elem.find(f'{{{NS_A}}}prstDash')
    if prstDash is not None:
        info["dash"] = prstDash.get('val')
    
    return info
# ---------------------------------------------------------------------------
# Master slide text styles extraction (9 levels)
# ---------------------------------------------------------------------------

def extract_text_styles_from_master(slide_master) -> Dict[str, Any]:
    """
    Extract text styles defined at master level (9 levels for body, title, other).
    These define the default formatting for each outline level.
    """
    styles = {
        "title_style": [],
        "body_style": [],
        "other_style": [],
    }
    
    try:
        if not hasattr(slide_master, '_element'):
            return styles
        
        elem = slide_master._element
        txStyles = elem.find(f'.//{{{NS_P}}}txStyles')
        
        if txStyles is None:
            warn("No txStyles found in slide master")
            return styles
        
        # Title style
        titleStyle = txStyles.find(f'{{{NS_P}}}titleStyle')
        if titleStyle is not None:
            styles["title_style"] = _extract_style_levels(titleStyle)
        
        # Body style
        bodyStyle = txStyles.find(f'{{{NS_P}}}bodyStyle')
        if bodyStyle is not None:
            styles["body_style"] = _extract_style_levels(bodyStyle)
        
        # Other style (for shapes that aren't title/body)
        otherStyle = txStyles.find(f'{{{NS_P}}}otherStyle')
        if otherStyle is not None:
            styles["other_style"] = _extract_style_levels(otherStyle)
            
    except Exception as e:
        warn(f"Text styles extraction error: {e}")
    
    return styles


def _extract_style_levels(style_elem) -> List[Dict[str, Any]]:
    """Extract all 9 levels of text style."""
    levels = []
    
    # Level tags: lvl1pPr through lvl9pPr
    for i in range(1, 10):
        level_tag = f'lvl{i}pPr'
        level_elem = style_elem.find(f'{{{NS_A}}}{level_tag}')
        
        if level_elem is not None:
            level_info = _extract_paragraph_properties(level_elem)
            level_info["level"] = i
            levels.append(level_info)
    
    # Also check defPPr (default paragraph properties)
    defPPr = style_elem.find(f'{{{NS_A}}}defPPr')
    if defPPr is not None:
        default_info = _extract_paragraph_properties(defPPr)
        default_info["level"] = 0
        default_info["is_default"] = True
        levels.insert(0, default_info)
    
    return levels


def _extract_paragraph_properties(pPr_elem) -> Dict[str, Any]:
    """Extract paragraph properties from XML element."""
    info = {}
    
    # Alignment
    algn = pPr_elem.get('algn')
    if algn:
        info["alignment"] = algn
    
    # Margins/indentation
    marL = pPr_elem.get('marL')
    if marL:
        info["margin_left_emu"] = int(marL)
        info["margin_left_in"] = emu_to_inches(int(marL))
    
    marR = pPr_elem.get('marR')
    if marR:
        info["margin_right_emu"] = int(marR)
    
    indent = pPr_elem.get('indent')
    if indent:
        info["indent_emu"] = int(indent)
        info["indent_in"] = emu_to_inches(int(indent))
    
    # Default tab size
    defTabSz = pPr_elem.get('defTabSz')
    if defTabSz:
        info["default_tab_size_emu"] = int(defTabSz)
    
    # RTL
    rtl = pPr_elem.get('rtl')
    if rtl:
        info["rtl"] = rtl == '1' or rtl == 'true'
    
    # Line spacing
    lnSpc = pPr_elem.find(f'{{{NS_A}}}lnSpc')
    if lnSpc is not None:
        spcPct = lnSpc.find(f'{{{NS_A}}}spcPct')
        if spcPct is not None:
            val = spcPct.get('val')
            if val:
                info["line_spacing_pct"] = int(val) / 1000
        spcPts = lnSpc.find(f'{{{NS_A}}}spcPts')
        if spcPts is not None:
            val = spcPts.get('val')
            if val:
                info["line_spacing_pt"] = int(val) / 100
    
    # Space before/after
    spcBef = pPr_elem.find(f'{{{NS_A}}}spcBef')
    if spcBef is not None:
        spcPts = spcBef.find(f'{{{NS_A}}}spcPts')
        if spcPts is not None:
            val = spcPts.get('val')
            if val:
                info["space_before_pt"] = int(val) / 100
    
    spcAft = pPr_elem.find(f'{{{NS_A}}}spcAft')
    if spcAft is not None:
        spcPts = spcAft.find(f'{{{NS_A}}}spcPts')
        if spcPts is not None:
            val = spcPts.get('val')
            if val:
                info["space_after_pt"] = int(val) / 100
    
    # Bullet
    bullet_info = {}
    
    buNone = pPr_elem.find(f'{{{NS_A}}}buNone')
    if buNone is not None:
        bullet_info["type"] = "none"
    
    buChar = pPr_elem.find(f'{{{NS_A}}}buChar')
    if buChar is not None:
        bullet_info["type"] = "char"
        bullet_info["char"] = buChar.get('char')
    
    buAutoNum = pPr_elem.find(f'{{{NS_A}}}buAutoNum')
    if buAutoNum is not None:
        bullet_info["type"] = "auto_number"
        bullet_info["scheme"] = buAutoNum.get('type')
        start = buAutoNum.get('startAt')
        if start:
            bullet_info["start_at"] = int(start)
    
    buFont = pPr_elem.find(f'{{{NS_A}}}buFont')
    if buFont is not None:
        bullet_info["font"] = buFont.get('typeface')
    
    buClr = pPr_elem.find(f'{{{NS_A}}}buClr')
    if buClr is not None:
        for color_child in buClr:
            color_info = _extract_xml_color_direct(color_child)
            if color_info:
                bullet_info["color"] = color_info
    
    buSzPct = pPr_elem.find(f'{{{NS_A}}}buSzPct')
    if buSzPct is not None:
        val = buSzPct.get('val')
        if val:
            bullet_info["size_pct"] = int(val) / 1000
    
    if bullet_info:
        info["bullet"] = bullet_info
    
    # Default run properties (font)
    defRPr = pPr_elem.find(f'{{{NS_A}}}defRPr')
    if defRPr is not None:
        font_info = _extract_run_properties(defRPr)
        if font_info:
            info["default_font"] = font_info
    
    return info


def _extract_run_properties(rPr_elem) -> Dict[str, Any]:
    """Extract run (character) properties from XML element."""
    info = {}
    
    # Font size
    sz = rPr_elem.get('sz')
    if sz:
        info["size_pt"] = int(sz) / 100
    
    # Bold
    b = rPr_elem.get('b')
    if b:
        info["bold"] = b == '1' or b == 'true'
    
    # Italic
    i = rPr_elem.get('i')
    if i:
        info["italic"] = i == '1' or i == 'true'
    
    # Underline
    u = rPr_elem.get('u')
    if u:
        info["underline"] = u
    
    # Strikethrough
    strike = rPr_elem.get('strike')
    if strike:
        info["strikethrough"] = strike
    
    # Caps
    cap = rPr_elem.get('cap')
    if cap:
        info["caps"] = cap
    
    # Baseline (sub/superscript)
    baseline = rPr_elem.get('baseline')
    if baseline:
        info["baseline"] = int(baseline)
    
    # Character spacing
    spc = rPr_elem.get('spc')
    if spc:
        info["char_spacing"] = int(spc)
    
    # Kerning
    kern = rPr_elem.get('kern')
    if kern:
        info["kerning"] = int(kern)
    
    # Color
    solidFill = rPr_elem.find(f'{{{NS_A}}}solidFill')
    if solidFill is not None:
        for color_child in solidFill:
            color_info = _extract_xml_color_direct(color_child)
            if color_info:
                info["color"] = color_info
    
    # Latin font
    latin = rPr_elem.find(f'{{{NS_A}}}latin')
    if latin is not None:
        info["latin_font"] = latin.get('typeface')
    
    # East Asian font
    ea = rPr_elem.find(f'{{{NS_A}}}ea')
    if ea is not None:
        info["ea_font"] = ea.get('typeface')
    
    # Complex script font
    cs = rPr_elem.find(f'{{{NS_A}}}cs')
    if cs is not None:
        info["cs_font"] = cs.get('typeface')
    
    return info


# ---------------------------------------------------------------------------
# Color map extraction
# ---------------------------------------------------------------------------

def extract_color_map(element) -> Dict[str, str]:
    """Extract color map from slide master or layout."""
    color_map = {}
    
    try:
        clrMap = None
        if hasattr(element, '_element'):
            clrMap = element._element.find(f'.//{{{NS_P}}}clrMap')
        
        if clrMap is not None:
            # All possible color map attributes
            map_attrs = ['bg1', 'bg2', 'tx1', 'tx2', 
                        'accent1', 'accent2', 'accent3', 'accent4', 'accent5', 'accent6',
                        'hlink', 'folHlink']
            
            for attr in map_attrs:
                val = clrMap.get(attr)
                if val:
                    color_map[attr] = val
    except Exception:
        pass
    
    return color_map


# ---------------------------------------------------------------------------
# Slide dimensions
# ---------------------------------------------------------------------------

def extract_slide_dimensions(prs) -> Dict[str, Any]:
    """Extract slide dimensions with aspect ratio detection."""
    try:
        w = prs.slide_width
        h = prs.slide_height
        ratio = w / h
        
        # Detect aspect ratio
        if abs(ratio - 16/9) < 0.01:
            aspect = "16:9"
        elif abs(ratio - 4/3) < 0.01:
            aspect = "4:3"
        elif abs(ratio - 16/10) < 0.01:
            aspect = "16:10"
        elif abs(ratio - 3/2) < 0.01:
            aspect = "3:2"
        elif ratio > 1.6:
            aspect = "widescreen"
        elif ratio < 1.0:
            aspect = "portrait"
        else:
            aspect = "custom"
        
        return {
            "width_emu": w,
            "height_emu": h,
            "width_in": emu_to_inches(w),
            "height_in": emu_to_inches(h),
            "width_pt": emu_to_pt(w),
            "height_pt": emu_to_pt(h),
            "aspect_ratio": aspect,
            "aspect_ratio_decimal": round(ratio, 4),
        }
    except Exception as e:
        warn(f"Slide dimensions extraction error: {e}")
        return {}


# ---------------------------------------------------------------------------
# Slide classification (comprehensive)
# ---------------------------------------------------------------------------

LAYOUT_CLASSIFICATION_RULES = {
    # Title/cover slides
    "title_slide": [
        "title slide", "cover", "front page", "full area", "titelfolie",
        "deckblatt", "title only"
    ],
    # Section dividers
    "section_divider": [
        "section", "divider", "chapter", "interstitial", "break",
        "zwischenfolie", "kapitel"
    ],
    # Closing slides
    "closing": [
        "end", "thank", "closing", "conclusion", "next step", "questions",
        "q&a", "contact", "schluss", "ende"
    ],
    # Blank layouts
    "blank": [
        "blank", "leer", "empty"
    ],
    # Comparison layouts
    "comparison": [
        "comparison", "two content", "vergleich"
    ],
    # Picture-focused layouts
    "picture": [
        "picture", "image", "photo", "bild"
    ],
    # Quote/statement layouts
    "quote": [
        "quote", "statement", "zitat"
    ],
    # Keynote/highlight
    "keynote": [
        "key note", "keynote", "highlight", "key message"
    ],
}


def classify_layout_name(layout_name: str) -> str:
    """Classify layout by name pattern matching."""
    name = layout_name.lower().strip()
    
    # Check specific patterns first
    for classification, patterns in LAYOUT_CLASSIFICATION_RULES.items():
        for pattern in patterns:
            if pattern in name:
                return classification
    
    # Content layouts: explicitly named "content" or "grid"
    if "content" in name or "grid" in name:
        return "content"
    
    # Check for subtitle pattern (indicates title slide variant)
    if name in ("title", "titel") and "section" not in name:
        return "title_slide"
    
    # Default to content for anything else
    return "content"


def classify_slide(slide) -> str:
    """Classify slide by its layout."""
    try:
        return classify_layout_name(slide.slide_layout.name)
    except Exception:
        return "content"


# ---------------------------------------------------------------------------
# Layout blueprint extraction (comprehensive)
# ---------------------------------------------------------------------------

def extract_layout_blueprint(layout, layout_idx: int) -> Dict[str, Any]:
    """
    Extract COMPLETE layout blueprint.
    """
    stat("layouts_extracted")
    
    blueprint = {
        "layout_index": layout_idx,
        "name": layout.name,
        "classification": classify_layout_name(layout.name),
        "background": {},
        "inherit_background": True,
        "show_master_shapes": True,
        "color_map_override": {},
        "placeholders": [],
        "decorative_shapes": [],
        "placeholder_summary": {},
    }
    
    # Background
    try:
        blueprint["background"] = extract_fill_info(layout.background.fill)
    except Exception:
        pass
    
    # Check if layout has background override or inherits from master
    try:
        if hasattr(layout, '_element'):
            cSld = layout._element.find(f'.//{{{NS_P}}}cSld')
            if cSld is not None:
                bg = cSld.find(f'{{{NS_P}}}bg')
                if bg is not None:
                    blueprint["inherit_background"] = False
    except Exception:
        pass
    
    # Show master shapes flag
    try:
        if hasattr(layout, '_element'):
            elem = layout._element
            show_master = elem.get('showMasterSp')
            if show_master is not None:
                blueprint["show_master_shapes"] = show_master != '0'
    except Exception:
        pass
    
    # Color map override
    try:
        color_map = extract_color_map(layout)
        if color_map:
            blueprint["color_map_override"] = color_map
    except Exception:
        pass
    
    # Extract all shapes with z-order
    placeholder_types_found = defaultdict(int)
    
    for z_idx, shape in enumerate(layout.shapes):
        shape_info = extract_shape_info(shape, include_text_samples=True, z_order=z_idx)
        
        if shape_info.get("placeholder"):
            blueprint["placeholders"].append(shape_info)
            ph_type = shape_info["placeholder"].get("type", "UNKNOWN")
            placeholder_types_found[ph_type] += 1
        else:
            blueprint["decorative_shapes"].append(shape_info)
    
    blueprint["placeholder_summary"] = dict(placeholder_types_found)
    blueprint["placeholder_count"] = len(blueprint["placeholders"])
    blueprint["decorative_shape_count"] = len(blueprint["decorative_shapes"])
    
    return blueprint


# ---------------------------------------------------------------------------
# Actual slide extraction (comprehensive)
# ---------------------------------------------------------------------------

def extract_actual_slide(slide, slide_idx: int) -> Dict[str, Any]:
    """Extract COMPLETE actual slide information."""
    stat("slides_extracted")
    
    slide_data = {
        "slide_index": slide_idx,
        "layout_used": "",
        "layout_index": None,
        "classification": "content",
        "background": {},
        "placeholders": [],
        "shapes": [],
        "has_notes": False,
        "notes_preview": None,
        "slide_number_visible": None,
        "date_visible": None,
        "footer_visible": None,
    }
    
    try:
        slide_data["layout_used"] = slide.slide_layout.name
        # Get layout index
        master = slide.slide_layout.slide_master
        for idx, layout in enumerate(master.slide_layouts):
            if layout.name == slide.slide_layout.name:
                slide_data["layout_index"] = idx
                break
        slide_data["classification"] = classify_slide(slide)
    except Exception:
        pass
    
    try:
        slide_data["background"] = extract_fill_info(slide.background.fill)
    except Exception:
        pass
    
    # Check visibility of footer elements
    try:
        if hasattr(slide, '_element'):
            elem = slide._element
            
            # Slide number visibility
            sldNum = elem.find(f'.//{{{NS_P}}}hf')
            if sldNum is not None:
                slide_data["slide_number_visible"] = sldNum.get('sldNum') != '0'
                slide_data["date_visible"] = sldNum.get('dt') != '0'
                slide_data["footer_visible"] = sldNum.get('ftr') != '0'
    except Exception:
        pass
    
    # Extract all shapes
    for z_idx, shape in enumerate(slide.shapes):
        shape_info = extract_shape_info(shape, include_text_samples=True, z_order=z_idx)
        
        if shape_info.get("placeholder"):
            slide_data["placeholders"].append(shape_info)
        else:
            slide_data["shapes"].append(shape_info)
    
    # Notes
    try:
        notes_slide = slide.notes_slide
        notes_text = notes_slide.notes_text_frame.text.strip()
        if notes_text:
            slide_data["has_notes"] = True
            slide_data["notes_preview"] = notes_text[:300]
    except Exception:
        pass
    
    return slide_data


# ---------------------------------------------------------------------------
# Notes master and Handout master extraction
# ---------------------------------------------------------------------------

def extract_notes_master(prs) -> Optional[Dict[str, Any]]:
    """Extract notes master layout information."""
    try:
        notes_master = prs.notes_master
        if notes_master is None:
            return None
        
        info = {
            "background": {},
            "placeholders": [],
            "shapes": [],
        }
        
        try:
            info["background"] = extract_fill_info(notes_master.background.fill)
        except Exception:
            pass
        
        for z_idx, shape in enumerate(notes_master.shapes):
            shape_info = extract_shape_info(shape, include_text_samples=True, z_order=z_idx)
            if shape_info.get("placeholder"):
                info["placeholders"].append(shape_info)
            else:
                info["shapes"].append(shape_info)
        
        return info
    except Exception:
        return None


def extract_handout_master(prs) -> Optional[Dict[str, Any]]:
    """Extract handout master layout information."""
    try:
        # python-pptx doesn't directly expose handout master, access via XML
        if not hasattr(prs, 'part') or not hasattr(prs.part, 'related_parts'):
            return None
        
        # Check for handout master relationship
        for rel in prs.part.rels.values():
            if 'handoutMaster' in rel.reltype:
                return {"note": "Handout master exists - extraction requires direct XML parsing"}
        
        return None
    except Exception:
        return None


# ---------------------------------------------------------------------------
# Main extraction orchestrator (comprehensive)
# ---------------------------------------------------------------------------

def extract_master_slide(pptx_path: str) -> Dict[str, Any]:
    """
    COMPREHENSIVE extraction of all master slide and layout information.
    """
    global _extraction_warnings, _extraction_stats
    _extraction_warnings = []
    _extraction_stats = defaultdict(int)
    
    prs = Presentation(pptx_path)

    result = {
        "source_file": str(Path(pptx_path).name),
        "extraction_version": "2.0.0",
        
        # Document-level
        "slide_dimensions": extract_slide_dimensions(prs),
        
        # Theme (complete)
        "theme": extract_theme_colors(prs),
        
        # Master slides (all of them)
        "slide_masters": [],
        
        # All layouts (complete blueprints)
        "all_layout_blueprints": [],
        
        # Canonical layouts (best per classification)
        "canonical_layouts": {},
        
        # Actual slides (organized by type)
        "actual_slides": {
            "title_slide": None,
            "section_dividers": [],
            "content_slides": [],
            "closing_slides": [],
            "blank_slides": [],
            "comparison_slides": [],
            "picture_slides": [],
            "quote_slides": [],
            "keynote_slides": [],
            "other": []
        },
        
        # Notes and handout masters
        "notes_master": None,
        "handout_master": None,
        
        # Synthesized summary
        "style_summary": {},
        
        # Extraction metadata
        "extraction_stats": {},
        "extraction_warnings": [],
        "completeness_check": {},
    }

    # ---- Extract all slide masters ----
    for master_idx, slide_master in enumerate(prs.slide_masters):
        stat("slide_masters")
        
        master_info = {
            "master_index": master_idx,
            "background": {},
            "color_map": {},
            "text_styles": {},
            "placeholders": [],
            "decorative_shapes": [],
            "placeholder_summary": {},
        }
        
        # Background
        try:
            master_info["background"] = extract_fill_info(slide_master.background.fill)
        except Exception as e:
            warn(f"Master {master_idx} background extraction error: {e}")
        
        # Color map
        try:
            master_info["color_map"] = extract_color_map(slide_master)
        except Exception:
            pass
        
        # Text styles (9 levels)
        try:
            master_info["text_styles"] = extract_text_styles_from_master(slide_master)
        except Exception as e:
            warn(f"Master {master_idx} text styles extraction error: {e}")
        
        # All shapes with z-order
        placeholder_types = defaultdict(int)
        for z_idx, shape in enumerate(slide_master.shapes):
            si = extract_shape_info(shape, include_text_samples=True, z_order=z_idx)
            if si.get("placeholder"):
                master_info["placeholders"].append(si)
                ph_type = si["placeholder"].get("type", "UNKNOWN")
                placeholder_types[ph_type] += 1
            else:
                master_info["decorative_shapes"].append(si)
        
        master_info["placeholder_summary"] = dict(placeholder_types)
        master_info["placeholder_count"] = len(master_info["placeholders"])
        master_info["decorative_shape_count"] = len(master_info["decorative_shapes"])
        
        result["slide_masters"].append(master_info)

    # ---- Extract all layout blueprints ----
    blueprints_by_class = defaultdict(list)
    
    if prs.slide_masters:
        for idx, layout in enumerate(prs.slide_masters[0].slide_layouts):
            bp = extract_layout_blueprint(layout, idx)
            result["all_layout_blueprints"].append(bp)
            blueprints_by_class[bp["classification"]].append(bp)

    # ---- Pick canonical (best) layout per classification ----
    for cls, bps in blueprints_by_class.items():
        if not bps:
            continue
            
        if cls == "content":
            # Prefer layouts named "Content | X" over "Grid | X" or others
            content_named = [bp for bp in bps if "content" in bp["name"].lower()]
            if content_named:
                # Further prefer simple content layouts
                simple = [bp for bp in content_named 
                         if bp["name"].lower() in ("content | 1", "content", "content only")]
                result["canonical_layouts"][cls] = simple[0] if simple else content_named[0]
            else:
                result["canonical_layouts"][cls] = bps[0]
        elif cls == "title_slide":
            # Prefer "Title Slide" over "Title Only"
            titled = [bp for bp in bps if "title slide" in bp["name"].lower()]
            result["canonical_layouts"][cls] = titled[0] if titled else bps[0]
        else:
            result["canonical_layouts"][cls] = bps[0]

    # ---- Extract actual slides ----
    slide_classification_map = {
        "title_slide": "title_slide",
        "section_divider": "section_dividers",
        "content": "content_slides",
        "closing": "closing_slides",
        "blank": "blank_slides",
        "comparison": "comparison_slides",
        "picture": "picture_slides",
        "quote": "quote_slides",
        "keynote": "keynote_slides",
    }
    
    for idx, slide in enumerate(prs.slides):
        sd = extract_actual_slide(slide, idx)
        cls = sd["classification"]
        
        if cls == "title_slide" and result["actual_slides"]["title_slide"] is None:
            result["actual_slides"]["title_slide"] = sd
        elif cls in slide_classification_map:
            key = slide_classification_map[cls]
            if isinstance(result["actual_slides"].get(key), list):
                result["actual_slides"][key].append(sd)
        else:
            result["actual_slides"]["other"].append(sd)

    # ---- Notes and handout masters ----
    result["notes_master"] = extract_notes_master(prs)
    result["handout_master"] = extract_handout_master(prs)

    # ---- Synthesize style summary ----
    result["style_summary"] = synthesize_style_summary(result)

    # ---- Extraction statistics and warnings ----
    result["extraction_stats"] = dict(_extraction_stats)
    result["extraction_warnings"] = _extraction_warnings
    
    # ---- Completeness check ----
    result["completeness_check"] = check_extraction_completeness(result)

    return result


def check_extraction_completeness(data: Dict[str, Any]) -> Dict[str, Any]:
    """
    Validate extraction completeness and report missing elements.
    """
    check = {
        "is_complete": True,
        "missing_elements": [],
        "warnings": [],
        "coverage": {},
    }
    
    # Check theme
    theme = data.get("theme", {})
    if not theme.get("colors"):
        check["missing_elements"].append("theme_colors")
        check["is_complete"] = False
    if not theme.get("font_scheme"):
        check["missing_elements"].append("font_scheme")
        check["is_complete"] = False
    
    # Check slide masters
    masters = data.get("slide_masters", [])
    if not masters:
        check["missing_elements"].append("slide_masters")
        check["is_complete"] = False
    else:
        master = masters[0]
        if not master.get("text_styles", {}).get("body_style"):
            check["warnings"].append("No body text styles found in master")
        if not master.get("text_styles", {}).get("title_style"):
            check["warnings"].append("No title text styles found in master")
    
    # Check layouts
    layouts = data.get("all_layout_blueprints", [])
    if not layouts:
        check["missing_elements"].append("layouts")
        check["is_complete"] = False
    else:
        check["coverage"]["layout_count"] = len(layouts)
        
        # Check for essential layout types
        classifications = set(l.get("classification") for l in layouts)
        if "content" not in classifications:
            check["warnings"].append("No content layout found")
        if "title_slide" not in classifications:
            check["warnings"].append("No title slide layout found")
    
    # Check actual slides
    actual = data.get("actual_slides", {})
    total_slides = 0
    for key, val in actual.items():
        if isinstance(val, list):
            total_slides += len(val)
        elif val is not None:
            total_slides += 1
    check["coverage"]["actual_slides"] = total_slides
    
    # Placeholder coverage
    all_placeholders = set()
    for layout in layouts:
        for ph in layout.get("placeholders", []):
            ph_type = ph.get("placeholder", {}).get("type")
            if ph_type:
                all_placeholders.add(ph_type)
    check["coverage"]["placeholder_types_found"] = list(all_placeholders)
    
    # Check for common elements
    if masters:
        master = masters[0]
        has_logo = any(s.get("is_image") for s in master.get("decorative_shapes", []))
        has_footer = any(s.get("is_footer_element") for s in master.get("decorative_shapes", []))
        check["coverage"]["has_logo_on_master"] = has_logo
        check["coverage"]["has_footer_elements"] = has_footer
    
    return check
# ---------------------------------------------------------------------------
# Style summary synthesis (comprehensive)
# ---------------------------------------------------------------------------

def synthesize_style_summary(data: Dict[str, Any]) -> Dict[str, Any]:
    """
    Synthesize a comprehensive, ready-to-use style summary.
    """
    summary = {
        "colors": {
            "palette": {},
            "semantic": {},
        },
        "typography": {
            "fonts_used": [],
            "heading_font": None,
            "body_font": None,
            "by_role": {},
            "text_styles_by_level": {},
        },
        "slide_dimensions": data.get("slide_dimensions", {}),
        "slide_templates": {},
        "master_decorative_elements": [],
        "footer_elements": [],
        "branding_elements": [],
        "design_notes": [],
        "layout_recommendations": {},
    }

    theme = data.get("theme", {})
    
    # ---- Colors ----
    colors = theme.get("colors", {})
    
    # Full palette
    summary["colors"]["palette"] = colors
    
    # Semantic color mapping
    color_semantics = {
        "dk1": "primary_text",
        "dk2": "secondary_text",
        "lt1": "background_primary",
        "lt2": "background_secondary",
        "accent1": "accent_primary",
        "accent2": "accent_secondary",
        "accent3": "accent_tertiary",
        "accent4": "accent_4",
        "accent5": "accent_5",
        "accent6": "accent_6",
        "hlink": "hyperlink",
        "folHlink": "followed_hyperlink",
    }
    
    for key, semantic_name in color_semantics.items():
        if key in colors:
            color_info = colors[key]
            if isinstance(color_info, dict):
                hex_val = color_info.get("hex")
            else:
                hex_val = color_info
            if hex_val:
                summary["colors"]["semantic"][semantic_name] = hex_val

    # ---- Typography ----
    fonts_seen = set()
    
    # From font scheme
    font_scheme = theme.get("font_scheme", {})
    
    major_font = font_scheme.get("majorFont", {})
    if major_font.get("latin"):
        summary["typography"]["heading_font"] = major_font["latin"]
        fonts_seen.add(major_font["latin"])
    
    minor_font = font_scheme.get("minorFont", {})
    if minor_font.get("latin"):
        summary["typography"]["body_font"] = minor_font["latin"]
        fonts_seen.add(minor_font["latin"])

    # Text styles from master
    masters = data.get("slide_masters", [])
    if masters:
        master = masters[0]
        text_styles = master.get("text_styles", {})
        
        # Title style summary
        title_styles = text_styles.get("title_style", [])
        if title_styles:
            lvl1 = next((s for s in title_styles if s.get("level") == 1), None)
            if lvl1:
                summary["typography"]["by_role"]["master_title"] = {
                    "font": lvl1.get("default_font", {}),
                    "bullet": lvl1.get("bullet"),
                    "alignment": lvl1.get("alignment"),
                }
                font_name = lvl1.get("default_font", {}).get("latin_font")
                if font_name:
                    fonts_seen.add(font_name)
        
        # Body style summary (all 9 levels)
        body_styles = text_styles.get("body_style", [])
        if body_styles:
            body_levels = {}
            for style in body_styles:
                level = style.get("level", 0)
                body_levels[f"level_{level}"] = {
                    "font": style.get("default_font", {}),
                    "bullet": style.get("bullet"),
                    "margin_left_in": style.get("margin_left_in"),
                    "indent_in": style.get("indent_in"),
                    "space_before_pt": style.get("space_before_pt"),
                    "space_after_pt": style.get("space_after_pt"),
                }
                font_name = style.get("default_font", {}).get("latin_font")
                if font_name:
                    fonts_seen.add(font_name)
            summary["typography"]["text_styles_by_level"]["body"] = body_levels
        
        # Collect fonts from decorative shapes
        for shape in master.get("decorative_shapes", []):
            _collect_fonts_from_shape(shape, fonts_seen)
        for shape in master.get("placeholders", []):
            _collect_fonts_from_shape(shape, fonts_seen)

    # Typography from actual slides
    actual = data.get("actual_slides", {})
    
    title_slide = actual.get("title_slide")
    if title_slide:
        summary["typography"]["by_role"]["title_slide_title"] = _first_ph_typography(
            title_slide, "TITLE"
        )
        summary["typography"]["by_role"]["title_slide_subtitle"] = _first_ph_typography(
            title_slide, "SUBTITLE"
        )
        _collect_fonts_from_slide(title_slide, fonts_seen)

    content_slides = actual.get("content_slides", [])
    if content_slides:
        cs = content_slides[0]
        summary["typography"]["by_role"]["content_slide_title"] = _first_ph_typography(cs, "TITLE")
        summary["typography"]["by_role"]["content_slide_body"] = _first_ph_typography(cs, "BODY")
        _collect_fonts_from_slide(cs, fonts_seen)

    section_dividers = actual.get("section_dividers", [])
    if section_dividers:
        sd = section_dividers[0]
        summary["typography"]["by_role"]["section_header"] = _first_ph_typography(sd, "TITLE")

    summary["typography"]["fonts_used"] = sorted(fonts_seen)

    # ---- Slide templates ----
    for cls, layout in data.get("canonical_layouts", {}).items():
        template = {
            "layout_name": layout.get("name"),
            "layout_index": layout.get("layout_index"),
            "classification": cls,
            "background": layout.get("background", {}),
            "placeholder_count": layout.get("placeholder_count", 0),
            "placeholder_summary": layout.get("placeholder_summary", {}),
            "placeholders": [],
        }
        
        for ph in layout.get("placeholders", []):
            pos = ph.get("position", {})
            tf = ph.get("text_frame", {})
            paras = tf.get("paragraphs", [])
            
            ph_info = {
                "placeholder_type": ph.get("placeholder", {}).get("type"),
                "placeholder_idx": ph.get("placeholder", {}).get("idx"),
                "category": ph.get("placeholder", {}).get("category"),
                "position_in": {
                    "left": pos.get("left_in"),
                    "top": pos.get("top_in"),
                    "width": pos.get("width_in"),
                    "height": pos.get("height_in"),
                },
                "position_emu": {
                    "left": pos.get("left_emu"),
                    "top": pos.get("top_emu"),
                    "width": pos.get("width_emu"),
                    "height": pos.get("height_emu"),
                },
                "text_frame": {
                    "margins_in": {k: v for k, v in tf.items() if "margin" in k},
                    "vertical_anchor": tf.get("vertical_anchor"),
                    "word_wrap": tf.get("word_wrap"),
                    "auto_size": tf.get("auto_size"),
                },
                "typography": paras[0].get("font", {}) if paras else {},
            }
            template["placeholders"].append(ph_info)
        
        summary["slide_templates"][cls] = template

    # ---- Master decorative elements ----
    if masters:
        master = masters[0]
        for shape in master.get("decorative_shapes", []):
            element = {
                "name": shape.get("name"),
                "shape_type": shape.get("shape_type"),
                "position": shape.get("position", {}),
                "z_order": shape.get("z_order"),
            }
            
            if shape.get("is_image"):
                element["is_logo_or_image"] = True
                element["image_info"] = shape.get("image_info")
                summary["branding_elements"].append(element)
            elif shape.get("is_footer_element"):
                element["footer_text"] = shape.get("footer_text")
                summary["footer_elements"].append(element)
            elif shape.get("is_confidential_label"):
                element["is_confidential"] = True
                summary["branding_elements"].append(element)
            else:
                element["fill"] = shape.get("fill")
                element["line"] = shape.get("line")
                element["effects"] = shape.get("effects")
                summary["master_decorative_elements"].append(element)

    # ---- Design notes ----
    dims = data.get("slide_dimensions", {})
    summary["design_notes"].append(
        f"Slide size: {dims.get('width_in')}\" × {dims.get('height_in')}\" "
        f"({dims.get('aspect_ratio')})"
    )
    
    if summary["typography"]["fonts_used"]:
        summary["design_notes"].append(
            f"Fonts: {', '.join(summary['typography']['fonts_used'])}"
        )
    
    if summary["typography"]["heading_font"]:
        summary["design_notes"].append(
            f"Theme heading font: {summary['typography']['heading_font']}"
        )
    
    if summary["typography"]["body_font"]:
        summary["design_notes"].append(
            f"Theme body font: {summary['typography']['body_font']}"
        )
    
    accent = summary["colors"]["semantic"].get("accent_primary")
    if accent:
        summary["design_notes"].append(f"Primary accent: {accent}")
    
    bg = summary["colors"]["semantic"].get("background_primary")
    if bg:
        summary["design_notes"].append(f"Primary background: {bg}")

    # Layout analysis
    layouts = data.get("all_layout_blueprints", [])
    layout_types = defaultdict(list)
    for layout in layouts:
        cls = layout.get("classification")
        layout_types[cls].append(layout.get("name"))
    
    summary["design_notes"].append(f"Layout types: {list(layout_types.keys())}")
    summary["design_notes"].append(f"Total layouts: {len(layouts)}")
    
    # Branding analysis
    if summary["branding_elements"]:
        logo_count = sum(1 for e in summary["branding_elements"] if e.get("is_logo_or_image"))
        summary["design_notes"].append(f"Branding elements on master: {len(summary['branding_elements'])}")
        if logo_count:
            summary["design_notes"].append(f"Logo/image elements: {logo_count}")
    
    if summary["footer_elements"]:
        footer_names = [e.get("name") for e in summary["footer_elements"]]
        summary["design_notes"].append(f"Footer elements: {footer_names}")

    # Picture placeholder analysis
    for cls, layout in data.get("canonical_layouts", {}).items():
        pic_phs = [p for p in layout.get("placeholders", []) if p.get("is_picture_placeholder")]
        if pic_phs:
            for pp in pic_phs:
                idx = pp.get("placeholder", {}).get("idx", "?")
                summary["design_notes"].append(
                    f"Layout '{layout.get('name')}' has PICTURE placeholder (idx={idx})"
                )

    # Layout recommendations
    for cls, layouts_list in layout_types.items():
        if len(layouts_list) > 1:
            summary["layout_recommendations"][cls] = {
                "available_layouts": layouts_list,
                "recommended": data.get("canonical_layouts", {}).get(cls, {}).get("name"),
            }

    return summary


def _collect_fonts_from_shape(shape: Dict, fonts_seen: set):
    """Collect font names from a shape."""
    tf = shape.get("text_frame", {})
    for para in tf.get("paragraphs", []):
        font = para.get("font", {})
        name = font.get("name")
        if name:
            fonts_seen.add(name)
        latin = font.get("latin_font")
        if latin:
            fonts_seen.add(latin)


def _collect_fonts_from_slide(slide: Dict, fonts_seen: set):
    """Collect font names from a slide."""
    for ph in slide.get("placeholders", []):
        _collect_fonts_from_shape(ph, fonts_seen)
    for shape in slide.get("shapes", []):
        _collect_fonts_from_shape(shape, fonts_seen)


def _first_ph_typography(slide_data: Dict, ph_type_keyword: str) -> Dict[str, Any]:
    """Get typography info from first matching placeholder."""
    if not slide_data:
        return {}
    
    for ph in slide_data.get("placeholders", []):
        ph_type = ph.get("placeholder", {}).get("type", "")
        if ph_type_keyword.upper() in ph_type.upper():
            tf = ph.get("text_frame", {})
            paras = tf.get("paragraphs", [])
            if paras:
                return {
                    "font": paras[0].get("font", {}),
                    "alignment": paras[0].get("alignment"),
                    "space_before_pt": paras[0].get("space_before_pt"),
                    "space_after_pt": paras[0].get("space_after_pt"),
                    "line_spacing": paras[0].get("line_spacing"),
                    "bullet": paras[0].get("bullet"),
                }
    return {}


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="COMPREHENSIVE PowerPoint master slide and layout style extractor",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python extract_master_styles.py presentation.pptx
  python extract_master_styles.py deck.pptx -o styles.json
  python extract_master_styles.py template.pptx --verbose

This tool extracts ALL styling information from PowerPoint files including:
  • Complete theme (colors, fonts, effects, format schemes)
  • Master slide with text styles for all 9 outline levels
  • All layout blueprints with full placeholder details
  • Decorative elements (logos, shapes, lines)
  • Footer and branding elements
  • Actual slide samples for reference
  • Comprehensive validation and completeness checks
        """
    )
    parser.add_argument("pptx_path", help="Path to the .pptx file")
    parser.add_argument("--output", "-o", help="Output JSON (default: <name>_styles.json)")
    parser.add_argument("--verbose", "-v", action="store_true", help="Show detailed extraction info")
    args = parser.parse_args()

    pptx_path = args.pptx_path
    output_path = args.output or (str(Path(pptx_path).stem) + "_styles.json")

    print(f"{'='*70}")
    print(f"COMPREHENSIVE POWERPOINT STYLE EXTRACTOR v2.0")
    print(f"{'='*70}")
    print(f"\nExtracting from: {pptx_path}")

    data = extract_master_slide(pptx_path)

    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, default=str, ensure_ascii=False)

    ss = data["style_summary"]
    stats = data.get("extraction_stats", {})
    warnings_list = data.get("extraction_warnings", [])
    completeness = data.get("completeness_check", {})

    print(f"\n✓ Extraction complete → {output_path}")
    
    print(f"\n{'='*70}")
    print("EXTRACTION STATISTICS")
    print(f"{'='*70}")
    print(f"  Slide masters:        {stats.get('slide_masters', 0)}")
    print(f"  Layouts extracted:    {stats.get('layouts_extracted', 0)}")
    print(f"  Slides extracted:     {stats.get('slides_extracted', 0)}")
    print(f"  Shapes processed:     {stats.get('shapes_extracted', 0)}")
    print(f"  Images found:         {stats.get('images', 0)}")
    print(f"  Tables found:         {stats.get('tables', 0)}")
    print(f"  Charts found:         {stats.get('charts', 0)}")
    print(f"  Groups found:         {stats.get('groups', 0)}")
    print(f"  SmartArt found:       {stats.get('smartart', 0)}")
    print(f"  Picture placeholders: {stats.get('picture_placeholders', 0)}")

    print(f"\n{'='*70}")
    print("STYLE SUMMARY")
    print(f"{'='*70}")
    
    dims = ss["slide_dimensions"]
    print(f"\n📐 Slide Dimensions:")
    print(f"   {dims.get('width_in')}\" × {dims.get('height_in')}\" ({dims.get('aspect_ratio')})")

    print(f"\n🎨 Color Palette:")
    for name, hex_val in ss["colors"]["semantic"].items():
        print(f"   {name:25s} {hex_val}")

    print(f"\n🔤 Typography:")
    print(f"   Fonts found: {', '.join(ss['typography']['fonts_used']) or 'none detected'}")
    if ss["typography"].get("heading_font"):
        print(f"   Theme heading: {ss['typography']['heading_font']}")
    if ss["typography"].get("body_font"):
        print(f"   Theme body: {ss['typography']['body_font']}")

    if args.verbose:
        print(f"\n   Typography by role:")
        for role, info in ss["typography"].get("by_role", {}).items():
            if info:
                font_info = info.get("font", {})
                font_name = font_info.get("name") or font_info.get("latin_font") or "inherited"
                font_size = font_info.get("size_pt", "?")
                print(f"     {role}: {font_name} @ {font_size}pt")

    print(f"\n📋 Slide Templates Extracted:")
    for cls, tmpl in ss["slide_templates"].items():
        ph_count = len(tmpl.get("placeholders", []))
        ph_summary = tmpl.get("placeholder_summary", {})
        print(f"   {cls:20s} '{tmpl['layout_name']}' — {ph_count} placeholder(s)")
        if args.verbose and ph_summary:
            ph_list = ", ".join(f"{k}:{v}" for k, v in ph_summary.items())
            print(f"                        Types: {ph_list}")

    print(f"\n🖼️  Master Elements:")
    print(f"   Decorative shapes: {len(ss['master_decorative_elements'])}")
    print(f"   Branding elements: {len(ss['branding_elements'])}")
    print(f"   Footer elements:   {len(ss['footer_elements'])}")

    if args.verbose and ss['branding_elements']:
        print(f"\n   Branding details:")
        for elem in ss['branding_elements']:
            print(f"     • {elem.get('name')}: {'Logo/Image' if elem.get('is_logo_or_image') else 'Other'}")

    print(f"\n📝 Design Notes:")
    for note in ss["design_notes"][:10]:  # Show first 10 notes
        print(f"   • {note}")
    if len(ss["design_notes"]) > 10:
        print(f"   ... and {len(ss['design_notes']) - 10} more notes")

    # Completeness check
    print(f"\n{'='*70}")
    print("COMPLETENESS CHECK")
    print(f"{'='*70}")
    
    if completeness.get("is_complete"):
        print("   ✓ Extraction complete - all major elements found")
    else:
        print("   ⚠ Some elements may be missing:")
        for elem in completeness.get("missing_elements", []):
            print(f"     • Missing: {elem}")
    
    if completeness.get("warnings"):
        print("\n   Warnings:")
        for warning in completeness["warnings"]:
            print(f"     • {warning}")
    
    coverage = completeness.get("coverage", {})
    if coverage:
        print(f"\n   Coverage:")
        print(f"     • Layouts: {coverage.get('layout_count', 0)}")
        print(f"     • Slides: {coverage.get('actual_slides', 0)}")
        print(f"     • Placeholder types: {len(coverage.get('placeholder_types_found', []))}")
        print(f"     • Logo on master: {'Yes' if coverage.get('has_logo_on_master') else 'No'}")
        print(f"     • Footer elements: {'Yes' if coverage.get('has_footer_elements') else 'No'}")

    # Extraction warnings
    if warnings_list:
        print(f"\n{'='*70}")
        print(f"EXTRACTION WARNINGS ({len(warnings_list)})")
        print(f"{'='*70}")
        for i, warning in enumerate(warnings_list[:5]):
            print(f"   {i+1}. {warning}")
        if len(warnings_list) > 5:
            print(f"   ... and {len(warnings_list) - 5} more warnings (see JSON for full list)")

    print(f"\n{'='*70}")
    print(f"Output saved to: {output_path}")
    print(f"{'='*70}\n")


if __name__ == "__main__":
    main()
